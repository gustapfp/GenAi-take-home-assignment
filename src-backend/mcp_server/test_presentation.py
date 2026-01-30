"""
Test file for the create_presentation function.
Run this file directly to test the presentation generation with sample data.

Usage:
    python -m mcp_server.test_presentation
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt

from core.consts import FILE_PATH
from core.logger_config import logger

DEFAULT_FONT_SIZE = Pt(18)
FONT_SIZE_WITH_IMAGE = Pt(14)

IMAGE_HEIGHT = Inches(2.5)
IMAGE_LEFT = Inches(5.0)
IMAGE_TOP = Inches(2.2)

BODY_WIDTH_WITH_IMAGE = Inches(4.5)


def create_presentation(filename: str, slides_content: str) -> str:
    """Create a PowerPoint presentation based on the given slides content.

    Args:
        filename (str): The filename of the presentation.
        slides_content (str): The slides content of the presentation.

    Returns:
        str: The message indicating the success or failure of the operation.
    """
    try:
        data = json.loads(slides_content)

        prs = Presentation()

        for slide_data in data:
            # -- Bullet layout --
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Check if this slide has an image
            image_path = slide_data.get("image")
            has_image = image_path and os.path.exists(image_path)

            # -- Title --
            title = slide.shapes.title
            if title:
                title.text = slide_data.get("title", "No Title")

            # -- Body --
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True

            # Adjust body width if image is present to prevent overlap
            if has_image:
                body_shape.width = BODY_WIDTH_WITH_IMAGE

            # Set font size based on whether image is present
            font_size = FONT_SIZE_WITH_IMAGE if has_image else DEFAULT_FONT_SIZE

            # Get the points to add
            points = slide_data.get("points", [])

            # Clear existing text and add bullet points
            if points:
                # Set the first point in the existing paragraph
                tf.paragraphs[0].text = points[0]
                tf.paragraphs[0].font.size = font_size
                tf.paragraphs[0].level = 0  # First level bullet

                # Add remaining points as new paragraphs
                for point in points[1:]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.font.size = font_size
                    p.level = 0  # First level bullet
                    p.space_before = Pt(6)  # Add spacing between points

            # -- Speaker Notes & Sources --
            speaker_notes = slide_data.get("speaker_notes", "")
            sources = slide_data.get("sources", [])
            if speaker_notes or sources:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                content = speaker_notes if speaker_notes else ""

                if sources:
                    if content:
                        content += "\n\n"
                    content += "Sources:\n" + "\n".join(f"- {url}" for url in sources)
                if text_frame:
                    text_frame.text = content

            # -- Image --
            if has_image:
                try:
                    slide.shapes.add_picture(
                        image_path,
                        left=IMAGE_LEFT,
                        top=IMAGE_TOP,
                        height=IMAGE_HEIGHT,
                    )
                except Exception as e:
                    logger.warning(f"Could not add image {image_path}: {e}")

        path = FILE_PATH / f"{filename}.pptx"
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path.resolve()))
        return f"Successfully saved presentation to {path}"
    except Exception as e:
        return f"Error creating PPT: {str(e)}"


# Sample test data
SAMPLE_SLIDES = [
    {
        "title": "Introduction to AI",
        "points": [
            "Artificial Intelligence is transforming industries worldwide",
            "Machine Learning enables computers to learn from data",
            "Deep Learning uses neural networks with multiple layers",
            "AI applications include healthcare, finance, and transportation",
        ],
        "speaker_notes": "Start with a brief overview of what AI means and its significance in today's world.",
        "sources": [
            "https://example.com/ai-overview",
            "https://example.com/ml-basics",
        ],
        "image": None,
    },
    {
        "title": "Key AI Technologies",
        "points": [
            "Natural Language Processing (NLP) - Understanding human language",
            "Computer Vision - Interpreting visual information",
            "Reinforcement Learning - Learning through trial and error",
            "Generative AI - Creating new content from learned patterns",
        ],
        "speaker_notes": "Explain each technology briefly and give real-world examples.",
        "sources": [
            "https://example.com/nlp-guide",
            "https://example.com/computer-vision",
        ],
        "image": None,
    },
    {
        "title": "AI in Business",
        "points": [
            "Automation of repetitive tasks saves time and reduces errors",
            "Predictive analytics helps in decision making",
            "Chatbots improve customer service efficiency",
            "Personalization increases customer engagement",
            "Cost reduction through optimized operations",
        ],
        "speaker_notes": "Focus on ROI and practical business benefits. Mention case studies if available.",
        "sources": [
            "https://example.com/ai-business-cases",
        ],
        "image": None,
    },
    {
        "title": "Future of AI",
        "points": [
            "AGI (Artificial General Intelligence) research continues",
            "Ethical AI and responsible development gaining importance",
            "AI regulation and governance frameworks emerging",
            "Integration with IoT and edge computing",
        ],
        "speaker_notes": "End with a forward-looking perspective and open questions for discussion.",
        "sources": [],
        "image": None,
    },
]


SAMPLE_CHART_PATH = str(FILE_PATH / "charts" / "chart_Amazon_Financial_Performance_1769752141.png")

SAMPLE_SLIDES_WITH_CHART = [
    {
        "title": "AI Market Growth",
        "points": [
            "The global AI market is experiencing rapid growth",
            "Expected to reach $500B by 2028",
            "North America leads in AI adoption",
            "Healthcare and finance are top sectors",
        ],
        "speaker_notes": "Reference the chart on the right showing market projections.",
        "sources": ["https://example.com/ai-market-report"],
        "image": SAMPLE_CHART_PATH,  # Test with actual chart
    },
    {
        "title": "Investment Trends",
        "points": [
            "VC funding in AI startups at all-time high",
            "Major tech companies investing billions in AI R&D",
            "Government initiatives supporting AI development",
        ],
        "speaker_notes": "Discuss the investment landscape and key players.",
        "sources": [],
        "image": None,
    },
]

SAMPLE_SLIDES_MIXED = [
    {
        "title": "Slide Without Image",
        "points": [
            "This slide demonstrates the default layout",
            "Text uses the standard font size (18pt)",
            "Body takes full width of the content area",
            "No image positioning adjustments needed",
        ],
        "speaker_notes": "Compare this with the next slide.",
        "sources": [],
        "image": None,
    },
    {
        "title": "Slide With Image",
        "points": [
            "This slide has a chart on the right",
            "Text uses smaller font size (14pt)",
            "Body width is reduced to avoid overlap",
            "Image is positioned on the right side",
            "Layout automatically adjusts for clarity",
        ],
        "speaker_notes": "Notice how the text and image don't overlap.",
        "sources": ["https://example.com/source"],
        "image": SAMPLE_CHART_PATH,
    },
]


if __name__ == "__main__":
    print("=" * 60)
    print("Testing create_presentation function")
    print("=" * 60)

    import os as os_check

    if os_check.path.exists(SAMPLE_CHART_PATH):
        print(f"\n✓ Sample chart found: {SAMPLE_CHART_PATH}")
    else:
        print(f"\n✗ Sample chart NOT found: {SAMPLE_CHART_PATH}")
        print("  Tests with images will skip the image embedding.")

    print("\n[Test 1] Creating basic presentation without images...")
    slides_json = json.dumps(SAMPLE_SLIDES)
    result = create_presentation("test_presentation_basic", slides_json)
    print(f"Result: {result}")

    print("\n[Test 2] Creating presentation with chart image...")
    slides_json = json.dumps(SAMPLE_SLIDES_WITH_CHART)
    result = create_presentation("test_presentation_with_chart", slides_json)
    print(f"Result: {result}")

    print("\n[Test 3] Creating presentation with mixed slides (image vs no image)...")
    slides_json = json.dumps(SAMPLE_SLIDES_MIXED)
    result = create_presentation("test_presentation_mixed", slides_json)
    print(f"Result: {result}")

    print("\n[Test 4] Creating empty presentation...")
    result = create_presentation("test_presentation_empty", json.dumps([]))
    print(f"Result: {result}")

    print("\n[Test 5] Creating single slide presentation...")
    single_slide = [
        {
            "title": "Single Slide Test",
            "points": ["Point 1", "Point 2", "Point 3"],
            "speaker_notes": "This is a single slide test.",
            "sources": [],
            "image": None,
        }
    ]
    result = create_presentation("test_presentation_single", json.dumps(single_slide))
    print(f"Result: {result}")

    print("\n" + "=" * 60)
    print(f"All presentations saved to: {FILE_PATH}")
    print("Open 'test_presentation_mixed.pptx' to compare slides with/without images.")
    print("=" * 60)
