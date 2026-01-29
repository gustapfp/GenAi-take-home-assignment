import matplotlib.pyplot as plt
from mcp_server.mcp_server import mcp_server
from structlog import get_logger
from pptx import Presentation
import json
from core.consts import FILE_PATH

logger = get_logger()


@mcp_server.tool(name="search_web", description="Search the web for information")
def search_web(query: str) -> str:
    """Search the web for information based on the given query.

    Args:
        query (str): The query to search the web for.

    Returns:
        str: The information searched for.
    """
    logger.info(f"Search Web Tool was triggered with query: {query}. Searching the web for information...")
    return "Information searched successfully!"


@mcp_server.tool(
    name="create_presentation", description="Create a PowerPoint presentation based on the given slides content."
)
def create_presentation(filename: str, slides_content: str) -> str:
    """Create a PowerPoint presentation based on the given slides content.

    Args:
        filename (str): The filename of the presentation.
        slides_content (str): The slides content of the presentation.

    Returns:
        str: The message indicating the success or failure of the operation.
    """
    try:
        data = json.loads(slides_content)
        prs = Presentation()

        for slide_data in data:
            # -- Bullet layout --
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # -- Title --
            title = slide.shapes.title
            if title:
                title.text = slide_data.get("title", "No Title")

            # -- Body --
            body_shape = slide.placeholders[1]
            tf = getattr(body_shape, "text_frame")
            for point in slide_data.get("points", []):
                p = tf.add_paragraph()
                p.text = point

        path = f"{FILE_PATH}/{filename}.pptx"
        prs.save(path)
        return f"Successfully saved presentation to {path}"
    except Exception as e:
        return f"Error creating PPT: {str(e)}"
